"""
Doküman Okuma Servisi
PDF, DOCX, PPTX ve TXT dosyalarından metin çıkarma fonksiyonları
"""

import os
import re
from typing import Optional, Tuple
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


class DocumentReader:
    """Çeşitli doküman formatlarından metin çıkarma sınıfı"""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """
        PDF dosyasından metin çıkarır
        
        Args:
            file_path: PDF dosyasının yolu
            
        Returns:
            Çıkarılan metin
            
        Raises:
            Exception: Dosya okuma hatası durumunda
        """
        try:
            reader = PdfReader(file_path)
            text = []
            
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
            
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"PDF okuma hatası: {str(e)}")
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """
        DOCX (Word) dosyasından metin çıkarır
        
        Args:
            file_path: DOCX dosyasının yolu
            
        Returns:
            Çıkarılan metin
            
        Raises:
            Exception: Dosya okuma hatası durumunda
        """
        try:
            doc = Document(file_path)
            text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            # Tablolardan da metin çıkar
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text.append(cell.text)
            
            return "\n".join(text)
        except Exception as e:
            raise Exception(f"DOCX okuma hatası: {str(e)}")
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """
        PPTX (PowerPoint) dosyasından metin çıkarır
        
        Args:
            file_path: PPTX dosyasının yolu
            
        Returns:
            Çıkarılan metin
            
        Raises:
            Exception: Dosya okuma hatası durumunda
        """
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if slide_text:
                    text.append(f"Slayt {slide_num}:\n" + "\n".join(slide_text))
            
            return "\n\n".join(text)
        except Exception as e:
            raise Exception(f"PPTX okuma hatası: {str(e)}")
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """
        TXT (Düz metin) dosyasından metin çıkarır
        
        Args:
            file_path: TXT dosyasının yolu
            
        Returns:
            Çıkarılan metin
            
        Raises:
            Exception: Dosya okuma hatası durumunda
        """
        try:
            # Farklı encoding'leri dene
            encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1254', 'iso-8859-9']
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            
            # Hiçbir encoding çalışmazsa hata fırlat
            raise Exception("Dosya encoding'i desteklenmiyor")
            
        except Exception as e:
            raise Exception(f"TXT okuma hatası: {str(e)}")
    
    @staticmethod
    def clean_text(text: str) -> str:
        """
        Metni temizler ve normalleştirir
        
        Args:
            text: Ham metin
            
        Returns:
            Temizlenmiş metin
        """
        # Çok fazla boşlukları temizle
        text = re.sub(r'\s+', ' ', text)
        
        # Çok fazla yeni satırları temizle
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        
        # Baştan ve sondan boşlukları temizle
        text = text.strip()
        
        return text
    
    @staticmethod
    def extract_text_from_file(file_path: str, file_extension: str) -> Tuple[str, Optional[str]]:
        """
        Dosya türüne göre uygun metodu çağırarak metin çıkarır
        
        Args:
            file_path: Dosyanın yolu
            file_extension: Dosya uzantısı (örn: 'pdf', 'docx')
            
        Returns:
            Tuple[str, Optional[str]]: (çıkarılan metin, hata mesajı)
            Başarılı ise (metin, None), başarısız ise ("", hata mesajı)
        """
        try:
            file_extension = file_extension.lower().strip('.')
            
            if file_extension == 'pdf':
                text = DocumentReader.extract_text_from_pdf(file_path)
            elif file_extension in ['docx', 'doc']:
                # DOC formatı için uyarı
                if file_extension == 'doc':
                    return ("", "Eski .doc formatı desteklenmiyor. Lütfen dosyayı .docx formatına dönüştürün.")
                text = DocumentReader.extract_text_from_docx(file_path)
            elif file_extension == 'pptx':
                text = DocumentReader.extract_text_from_pptx(file_path)
            elif file_extension == 'txt':
                text = DocumentReader.extract_text_from_txt(file_path)
            else:
                return ("", f"Desteklenmeyen dosya formatı: .{file_extension}")
            
            # Metni temizle
            text = DocumentReader.clean_text(text)
            
            # Metin boş mu kontrol et
            if not text or len(text.strip()) < 50:
                return ("", "Dosyadan yeterli metin çıkarılamadı. Dosya boş olabilir veya sadece resimlerden oluşuyor olabilir.")
            
            return (text, None)
            
        except Exception as e:
            return ("", f"Dosya okuma hatası: {str(e)}")
    
    @staticmethod
    def truncate_text(text: str, max_tokens: int = 12000) -> str:
        """
        Metni token limitine göre kısaltır
        Yaklaşık olarak 1 token = 4 karakter varsayımı
        
        Args:
            text: Kısaltılacak metin
            max_tokens: Maksimum token sayısı
            
        Returns:
            Kısaltılmış metin
        """
        max_chars = max_tokens * 4
        
        if len(text) <= max_chars:
            return text
        
        # Metni cümle sonlarından böl
        truncated = text[:max_chars]
        
        # Son cümleyi tamamla
        last_period = truncated.rfind('.')
        last_newline = truncated.rfind('\n')
        
        cutoff = max(last_period, last_newline)
        if cutoff > max_chars * 0.8:  # En az %80'ini al
            truncated = truncated[:cutoff + 1]
        
        return truncated + "\n\n[Not: Metin çok uzun olduğu için kısaltılmıştır]"

