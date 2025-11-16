"""
Pytest fixtures for StudyBuddy test suite
Merkezi fixture tanımları
"""

import pytest
import os
import tempfile
from datetime import datetime, timedelta
from unittest.mock import patch, MagicMock
from flask import Flask
from app import app as flask_app
from models import db, User, Document, Result
from config import Config


@pytest.fixture
def app():
    """Test Flask uygulaması oluştur"""
    # Test konfigürasyonu
    flask_app.config['TESTING'] = True
    # SQLite in-memory database'de her connection farklı bir database görür
    # Bu yüzden geçici bir dosya kullanarak tüm connection'ların aynı database'i görmesini sağlıyoruz
    test_db_fd, test_db_path = tempfile.mkstemp(suffix='.db')
    flask_app.config['SQLALCHEMY_DATABASE_URI'] = f'sqlite:///{test_db_path}'
    flask_app.config['WTF_CSRF_ENABLED'] = False
    flask_app.config['RATELIMIT_ENABLED'] = False
    flask_app.config['DEMO_MODE'] = True
    flask_app.config['SECRET_KEY'] = 'test-secret-key'
    flask_app.config['UPLOAD_FOLDER'] = tempfile.mkdtemp()
    
    # App context içinde çalıştır
    with flask_app.app_context():
        db.create_all()
        yield flask_app
        db.session.remove()
        db.drop_all()
    
    # Cleanup
    import shutil
    os.close(test_db_fd)
    if os.path.exists(test_db_path):
        os.unlink(test_db_path)
    if os.path.exists(flask_app.config['UPLOAD_FOLDER']):
        shutil.rmtree(flask_app.config['UPLOAD_FOLDER'])


@pytest.fixture
def client(app):
    """Flask test client oluştur"""
    return app.test_client()


@pytest.fixture
def db_session(app):
    """Test database session"""
    with app.app_context():
        yield db.session


@pytest.fixture
def user(db_session):
    """Test kullanıcısı oluştur"""
    test_user = User(
        email='test@example.com',
        username='testuser',
        subscription_plan='free',
        tokens_remaining=10,
        trial_ends_at=datetime.utcnow() + timedelta(days=7)
    )
    test_user.set_password('Test123!')
    db_session.add(test_user)
    db_session.commit()
    return test_user


@pytest.fixture
def premium_user(db_session):
    """Premium plan kullanıcısı oluştur"""
    premium = User(
        email='premium@example.com',
        username='premiumuser',
        subscription_plan='premium',
        tokens_remaining=60,
        trial_ends_at=datetime.utcnow() + timedelta(days=7)
    )
    premium.set_password('Test123!')
    db_session.add(premium)
    db_session.commit()
    return premium


@pytest.fixture
def authenticated_client(client, user):
    """Giriş yapmış client oluştur"""
    with client.session_transaction() as sess:
        sess['_user_id'] = str(user.id)
        sess['_fresh'] = True
    return client


@pytest.fixture
def sample_document(db_session, user):
    """Test Document kaydı oluştur"""
    doc = Document(
        file_hash='test_hash_123456789012345678901234567890',
        original_filename='test.pdf',
        file_type='pdf',
        file_size=1024,
        user_level='high_school',
        user_type='student',
        user_id=user.id
    )
    db_session.add(doc)
    db_session.commit()
    return doc


@pytest.fixture
def sample_result(db_session, sample_document):
    """Test Result kaydı oluştur"""
    import json
    
    result = Result(
        document_id=sample_document.id,
        summary='Test özet metni',
        multiple_choice=json.dumps([
            {
                'question': 'Test sorusu 1?',
                'options': ['A', 'B', 'C', 'D'],
                'correct_answer': 0,
                'explanation': 'Test açıklama'
            }
        ], ensure_ascii=False),
        short_answer=json.dumps([
            {
                'question': 'Test kısa cevap sorusu?',
                'answer': 'Test cevap'
            }
        ], ensure_ascii=False),
        fill_blank=json.dumps([
            {
                'question': 'Test boş doldurma _____ sorusu.',
                'answer': 'test',
                'options': ['test', 'test2', 'test3', 'test4']
            }
        ], ensure_ascii=False),
        true_false=json.dumps([
            {
                'statement': 'Test ifadesi',
                'is_true': True,
                'explanation': 'Test açıklama'
            }
        ], ensure_ascii=False),
        flashcards=json.dumps([
            {
                'front': 'Test ön yüz',
                'back': 'Test arka yüz'
            }
        ], ensure_ascii=False),
        ai_model='gpt-3.5-turbo',
        token_used=100,
        processing_time=2.5
    )
    db_session.add(result)
    db_session.commit()
    return result


@pytest.fixture
def mock_openai():
    """OpenAI API mock'u"""
    with patch('openai.OpenAI') as mock_client:
        mock_instance = MagicMock()
        mock_response = MagicMock()
        mock_response.choices = [MagicMock()]
        mock_response.choices[0].message.content = 'Mock AI response'
        mock_instance.chat.completions.create.return_value = mock_response
        mock_client.return_value = mock_instance
        yield mock_instance


@pytest.fixture
def sample_pdf_path(tmp_path):
    """Test PDF dosyasının yolunu döndür"""
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        
        pdf_path = tmp_path / 'sample.pdf'
        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        c.drawString(100, 750, "Test PDF Content")
        c.drawString(100, 700, "This is a test PDF file for testing purposes.")
        c.save()
        return str(pdf_path)
    except ImportError:
        # ReportLab yoksa basit bir PDF oluştur
        pdf_path = tmp_path / 'sample.pdf'
        with open(pdf_path, 'wb') as f:
            # Minimal PDF header
            f.write(b'%PDF-1.4\n')
            f.write(b'1 0 obj\n<< /Type /Catalog >>\nendobj\n')
            f.write(b'xref\n0 1\ntrailer\n<< /Root 1 0 R >>\nstartxref\n0\n%%EOF')
        return str(pdf_path)


@pytest.fixture
def sample_docx_path(tmp_path):
    """Test DOCX dosyasının yolunu döndür"""
    try:
        from docx import Document
        
        docx_path = tmp_path / 'sample.docx'
        doc = Document()
        doc.add_paragraph('Test DOCX Content')
        doc.add_paragraph('This is a test DOCX file for testing purposes.')
        doc.save(str(docx_path))
        return str(docx_path)
    except ImportError:
        # python-docx yoksa basit bir ZIP (DOCX formatı) oluştur
        import zipfile
        docx_path = tmp_path / 'sample.docx'
        with zipfile.ZipFile(docx_path, 'w') as z:
            z.writestr('word/document.xml', '<?xml version="1.0"?><document><body><p>Test content</p></body></document>')
        return str(docx_path)


@pytest.fixture
def sample_txt_path(tmp_path):
    """Test TXT dosyasının yolunu döndür"""
    txt_path = tmp_path / 'sample.txt'
    with open(txt_path, 'w', encoding='utf-8') as f:
        f.write('Test TXT Content\n')
        f.write('This is a test text file for testing purposes.\n')
        f.write('It contains multiple lines of text.')
    return str(txt_path)


@pytest.fixture
def sample_pptx_path(tmp_path):
    """Test PPTX dosyasının yolunu döndür"""
    try:
        from pptx import Presentation
        
        pptx_path = tmp_path / 'sample.pptx'
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Test Slide"
        content = slide.placeholders[1]
        content.text = "This is a test PowerPoint slide."
        prs.save(str(pptx_path))
        return str(pptx_path)
    except ImportError:
        # python-pptx yoksa basit bir ZIP (PPTX formatı) oluştur
        import zipfile
        pptx_path = tmp_path / 'sample.pptx'
        with zipfile.ZipFile(pptx_path, 'w') as z:
            z.writestr('ppt/presentation.xml', '<?xml version="1.0"?><presentation><slide><text>Test content</text></slide></presentation>')
        return str(pptx_path)


@pytest.fixture
def demo_mode_true(monkeypatch):
    """Config.DEMO_MODE'u True yapar ve test sonunda geri yükler"""
    original_value = Config.DEMO_MODE
    monkeypatch.setattr('config.Config.DEMO_MODE', True)
    yield
    monkeypatch.setattr('config.Config.DEMO_MODE', original_value)


@pytest.fixture
def demo_mode_false(monkeypatch):
    """Config.DEMO_MODE'u False yapar ve test sonunda geri yükler"""
    original_value = Config.DEMO_MODE
    monkeypatch.setattr('config.Config.DEMO_MODE', False)
    yield
    monkeypatch.setattr('config.Config.DEMO_MODE', original_value)


@pytest.fixture
def completed_payment(db_session, user):
    """Completed status'ta payment record oluştur"""
    from models import Payment
    payment = Payment(
        user_id=user.id,
        amount=49.99,
        currency='TRY',
        status='completed',
        stripe_session_id='cs_test_completed_123',
        stripe_payment_intent_id='pi_test_123',
        plan_type='premium',
        billing_period='monthly',
        invoice_number='INV-2025-00001',
        completed_at=datetime.utcnow()
    )
    db_session.add(payment)
    db_session.commit()
    return payment


@pytest.fixture
def standard_user(db_session):
    """Standard plan kullanıcısı oluştur"""
    from models import User
    standard = User(
        email='standard@example.com',
        username='standarduser',
        subscription_plan='standard',
        tokens_remaining=25,
        trial_ends_at=datetime.utcnow() + timedelta(days=7)
    )
    standard.set_password('Test123!')
    db_session.add(standard)
    db_session.commit()
    return standard


